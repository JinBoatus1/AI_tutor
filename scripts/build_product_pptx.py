"""Generate PRODUCT.pptx — run: python scripts/build_product_pptx.py from repo root."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    out = root / "PRODUCT.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "AI Tutor",
        "Textbook-first, syllabus-aware learning companion (FOCS)",
    )

    add_bullet_slide(
        prs,
        "What It Is",
        [
            "Web app: chat tutor + progress map + auto grader",
            "Grounded in the official FOCS textbook (PDF + structured outline)",
            "Designed for trustworthy, exam-aligned help—not generic web answers",
        ],
    )

    add_bullet_slide(
        prs,
        "Learning Mode",
        [
            "Natural-language Q&A; images, screenshots, paste, PDF (pages rendered server-side)",
            "Matches topics to FOCS; shows real textbook pages in a side panel (full section range)",
            "Short intake: new vs review, progress, target sections",
            "Confidence only on problem-solving turns (not navigation/planning)",
            "Solves: hallucinations & off-syllabus explanations",
        ],
    )

    add_bullet_slide(
        prs,
        "FOCS Tree & Textbook",
        [
            "FOCS.json: chapters, sections, page ranges",
            "Single source of truth with the FOCS PDF",
            "Solves: wrong pages / misleading crops on broad topics (e.g. Induction)",
        ],
    )

    add_bullet_slide(
        prs,
        "My Learning Bar",
        [
            "Full FOCS outline; teal = learned, gray = not yet",
            "Click to toggle; same student ID as Learning Mode",
            "Solves: invisible, scattered progress across chats",
        ],
    )

    add_bullet_slide(
        prs,
        "Hidden Progress Bar (Backend)",
        [
            "Per-student JSON: current, learned, planned, confusion counts",
            "Heuristics: chapter / subsection progress; prompt constraints for the tutor",
            "Solves: teaching too far ahead; remediation when “learned” still feels hard",
        ],
    )

    add_bullet_slide(
        prs,
        "Auto Grader",
        [
            "Custom grading prompt + student text + optional images",
            "Solves: fast draft feedback when no human grader is available",
        ],
    )

    add_bullet_slide(
        prs,
        "Session Memory",
        [
            "Per FOCS subtopic: summaries + optional full Q&A via tool",
            "Solves: continuity across long or multi-session threads",
        ],
    )

    add_bullet_slide(
        prs,
        "Problems → Solutions (Summary)",
        [
            "Trust → PDF + FOCS + page images",
            "Structure → Learning bar + intake + hidden bar",
            "Scope → bar-informed prompting",
            "Multimodal → images + PDF rendering",
            "Feedback → Auto Grader",
            "Continuity → subtopic memory",
        ],
    )

    add_bullet_slide(
        prs,
        "Stack (High Level)",
        [
            "Frontend: React, Vite",
            "Backend: FastAPI, PyMuPDF, OpenAI-compatible API",
            "Student ID in localStorage → one JSON bar file per learner",
        ],
    )

    add_bullet_slide(
        prs,
        "One-Liner",
        [
            "A textbook-first AI study companion: answers in the language of FOCS,",
            "shows real pages, respects what you marked as learned,",
            "and supports images, PDFs, grading help, and memory on demand.",
        ],
    )

    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
